#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX Processor - PowerPoint Dosyalarını Analiz ve İşleme Modülü

Bu modül herhangi bir PPTX dosyasını okuyarak:
1. Her sayfadaki metinleri ayırır ve metin dosyası olarak kaydeder
2. Her sayfadaki görselleri (resimler, şekiller) çıkarır ve kaydeder
3. Her sayfanın yapısını analiz eder ve rapor oluşturur

Kullanım:
    python pptx_processor.py <pptx_dosyasi.pptx>
    
    veya Python içinde:
    from pptx_processor import PPTXProcessor
    processor = PPTXProcessor("dosya.pptx")
    processor.process()
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PPTXProcessor:
    """PowerPoint dosyalarını analiz ve işleme sınıfı"""
    
    def __init__(self, pptx_path):
        """
        Args:
            pptx_path (str): İşlenecek PPTX dosyasının yolu
        """
        self.pptx_path = Path(pptx_path)
        if not self.pptx_path.exists():
            raise FileNotFoundError(f"PPTX dosyası bulunamadı: {pptx_path}")
        
        self.prs = Presentation(str(self.pptx_path))
        self.output_dir = self.pptx_path.parent / f"{self.pptx_path.stem}_output"
        self.output_dir.mkdir(exist_ok=True)
        
        # Alt dizinler
        self.text_dir = self.output_dir / "texts"
        self.image_dir = self.output_dir / "images"
        self.text_dir.mkdir(exist_ok=True)
        self.image_dir.mkdir(exist_ok=True)
        
        self.report = []
        
    def extract_text_from_shape(self, shape):
        """Bir şekilden metin çıkarır"""
        if hasattr(shape, "text") and shape.text.strip():
            return shape.text.strip()
        elif hasattr(shape, "text_frame"):
            text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                para_text = ""
                for run in paragraph.runs:
                    para_text += run.text
                if para_text.strip():
                    text_parts.append(para_text.strip())
            return "\n".join(text_parts) if text_parts else None
        return None
    
    def extract_image_from_shape(self, shape, slide_num, img_num):
        """Bir şekilden görsel çıkarır ve kaydeder"""
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                
                # Görsel formatını belirle
                ext = image.ext
                if not ext:
                    ext = "png"
                
                # Dosya adı oluştur
                filename = f"slide_{slide_num:02d}_image_{img_num:02d}.{ext}"
                filepath = self.image_dir / filename
                
                # Görseli kaydet
                with open(filepath, 'wb') as f:
                    f.write(image_bytes)
                
                return filename
        except Exception as e:
            self.report.append(f"   ⚠️  Görsel çıkarma hatası: {str(e)}")
        
        return None
    
    def process_slide(self, slide, slide_num):
        """Tek bir slaytı işler"""
        print(f"\n{'=' * 60}")
        print(f"SAYFA {slide_num} İŞLENİYOR...")
        print('=' * 60)
        
        slide_report = [f"\n## SAYFA {slide_num}"]
        
        # Metin ve görselleri topla
        texts = []
        images = []
        img_counter = 1
        
        for shape in slide.shapes:
            # Metin çıkar
            text = self.extract_text_from_shape(shape)
            if text:
                texts.append(text)
                print(f"✓ Metin bulundu: {text[:50]}...")
            
            # Görsel çıkar
            img_filename = self.extract_image_from_shape(shape, slide_num, img_counter)
            if img_filename:
                images.append(img_filename)
                print(f"✓ Görsel kaydedildi: {img_filename}")
                img_counter += 1
        
        # Metinleri dosyaya kaydet
        if texts:
            text_filename = f"slide_{slide_num:02d}_text.txt"
            text_filepath = self.text_dir / text_filename
            with open(text_filepath, 'w', encoding='utf-8') as f:
                f.write(f"# SAYFA {slide_num}\n\n")
                for i, text in enumerate(texts, 1):
                    f.write(f"## Metin Bloğu {i}\n")
                    f.write(f"{text}\n\n")
            
            slide_report.append(f"- **Metin Sayısı:** {len(texts)}")
            slide_report.append(f"- **Metin Dosyası:** `{text_filename}`")
            print(f"✓ Metinler kaydedildi: {text_filename}")
        else:
            slide_report.append("- **Metin:** Yok")
            print("⚠️  Bu sayfada metin bulunamadı")
        
        # Görseller
        if images:
            slide_report.append(f"- **Görsel Sayısı:** {len(images)}")
            slide_report.append(f"- **Görsel Dosyaları:**")
            for img in images:
                slide_report.append(f"  - `{img}`")
        else:
            slide_report.append("- **Görsel:** Yok")
            print("⚠️  Bu sayfada görsel bulunamadı")
        
        self.report.extend(slide_report)
        
    def process(self):
        """Tüm PPTX dosyasını işler"""
        print('\n' + '#' * 60)
        print(f"PPTX İŞLEME BAŞLADI: {self.pptx_path.name}")
        print('#' * 60)
        print(f"Toplam Sayfa Sayısı: {len(self.prs.slides)}")
        print(f"Çıktı Dizini: {self.output_dir}")
        
        self.report.append(f"# PPTX İŞLEME RAPORU")
        self.report.append(f"\n**Dosya:** `{self.pptx_path.name}`")
        self.report.append(f"**Toplam Sayfa:** {len(self.prs.slides)}")
        self.report.append(f"**İşlem Tarihi:** {Path(__file__).stat().st_mtime}")
        
        # Her sayfayı işle
        for slide_num, slide in enumerate(self.prs.slides, 1):
            self.process_slide(slide, slide_num)
        
        # Raporu kaydet
        self.save_report()
        
        print('\n' + '#' * 60)
        print("İŞLEM TAMAMLANDI!")
        print('#' * 60)
        print(f"✓ Tüm sayfalar işlendi: {len(self.prs.slides)} sayfa")
        print(f"✓ Metinler kaydedildi: {self.text_dir}")
        print(f"✓ Görseller kaydedildi: {self.image_dir}")
        print(f"✓ Rapor oluşturuldu: {self.output_dir / 'RAPOR.md'}")
        
    def save_report(self):
        """İşlem raporunu kaydeder"""
        report_path = self.output_dir / "RAPOR.md"
        with open(report_path, 'w', encoding='utf-8') as f:
            f.write("\n".join(self.report))
        print(f"\n✓ Rapor kaydedildi: {report_path}")
    
    def verify_all_pages(self):
        """Tüm sayfaların doğru şekilde işlendiğini kontrol eder"""
        print('\n' + '=' * 60)
        print("DOĞRULAMA YAPILIYOR...")
        print('=' * 60)
        
        total_slides = len(self.prs.slides)
        processed_texts = len(list(self.text_dir.glob("slide_*.txt")))
        processed_images = len(list(self.image_dir.glob("slide_*.png"))) + \
                          len(list(self.image_dir.glob("slide_*.jpg"))) + \
                          len(list(self.image_dir.glob("slide_*.jpeg")))
        
        print(f"\n📊 DOĞRULAMA SONUÇLARI:")
        print(f"   - Toplam Sayfa: {total_slides}")
        print(f"   - İşlenen Metin Dosyası: {processed_texts}")
        print(f"   - İşlenen Görsel: {processed_images}")
        
        if processed_texts > 0:
            print(f"   ✅ Metinler başarıyla işlendi")
        else:
            print(f"   ⚠️  Hiç metin bulunamadı veya işlenemedi")
        
        if processed_images > 0:
            print(f"   ✅ Görseller başarıyla işlendi")
        else:
            print(f"   ⚠️  Hiç görsel bulunamadı veya işlenemedi")
        
        # Detaylı sayfa kontrolü
        print(f"\n📄 SAYFA DETAYLARI:")
        for i in range(1, total_slides + 1):
            text_file = self.text_dir / f"slide_{i:02d}_text.txt"
            has_text = "✓" if text_file.exists() else "✗"
            
            image_files = list(self.image_dir.glob(f"slide_{i:02d}_image_*"))
            img_count = len(image_files)
            
            print(f"   Sayfa {i:2d}: Metin [{has_text}]  Görsel [{img_count} adet]")
        
        print('\n' + '=' * 60)
        print("DOĞRULAMA TAMAMLANDI")
        print('=' * 60 + '\n')


def main():
    """Komut satırı arayüzü"""
    if len(sys.argv) < 2:
        print("Kullanım: python pptx_processor.py <pptx_dosyasi.pptx>")
        print("\nÖrnek:")
        print("  python pptx_processor.py sunum.pptx")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    
    try:
        processor = PPTXProcessor(pptx_file)
        processor.process()
        processor.verify_all_pages()
    except Exception as e:
        print(f"\n❌ HATA: {str(e)}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
