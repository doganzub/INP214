#!/usr/bin/env python3
"""
PPTX İşleme Modülü (PPTX Processor)
====================================
Herhangi bir PPTX dosyasını analiz eder:
  - Yazı kısımlarını metin olarak çıkarır (Markdown formatında)
  - Görselleri görsel dosyası olarak kaydeder
  - Her slaytı ayrı ayrı işler
  - İşlem sonrası tüm sayfaları doğrulama kontrolü yapar

Kullanım:
    python pptx_processor.py dosya.pptx [--output-dir çıktı_klasörü]

Çıktı Yapısı:
    {dosya_adı}_output/
    ├── report.md              # Genel rapor
    ├── slide_01/
    │   ├── content.md         # Slaytın metin içeriği
    │   ├── image_01.png       # Çıkarılan görseller
    │   └── image_02.jpg
    ├── slide_02/
    │   ├── content.md
    │   └── image_01.png
    └── ...
"""

import sys
import argparse
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("HATA: python-pptx kütüphanesi gerekli!")
    print("Kurulum: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
    import io
    HAS_PIL = True
except ImportError:
    HAS_PIL = False


# ─────────────────────────────────────────────
# Ana İşlem Fonksiyonu
# ─────────────────────────────────────────────

def process_pptx(pptx_path, output_dir=None):
    """
    PPTX dosyasını işler: metin ve görselleri ayırarak çıkarır.

    Args:
        pptx_path: PPTX dosyasının yolu
        output_dir: Çıktı klasörü (varsayılan: {dosya_adı}_output)

    Returns:
        dict: İşlem raporu
    """
    pptx_path = Path(pptx_path).resolve()
    if not pptx_path.exists():
        raise FileNotFoundError(f"Dosya bulunamadı: {pptx_path}")

    if output_dir is None:
        output_dir = pptx_path.parent / f"{pptx_path.stem}_output"
    else:
        output_dir = Path(output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))

    # Slayt boyutlarını cm'ye çevir
    emu_to_cm = 914400 / 2.54
    w_cm = prs.slide_width / emu_to_cm
    h_cm = prs.slide_height / emu_to_cm

    report = {
        "dosya": str(pptx_path),
        "cikti_klasoru": str(output_dir),
        "slayt_sayisi": len(prs.slides),
        "slayt_boyutu": f"{w_cm:.1f} cm x {h_cm:.1f} cm",
        "slaytlar": [],
        "toplam_metin": 0,
        "toplam_gorsel": 0,
        "hatalar": [],
    }

    print(f"PPTX İşleniyor: {pptx_path.name}")
    print(f"Slayt sayısı: {len(prs.slides)}")
    print(f"Slayt boyutu: {report['slayt_boyutu']}")
    print(f"{'─'*60}")

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        slide_dir = output_dir / f"slide_{slide_num:02d}"
        slide_dir.mkdir(parents=True, exist_ok=True)

        print(f"\n📄 Slayt {slide_num}/{len(prs.slides)} işleniyor...")
        slide_info = _process_slide(slide, slide_num, slide_dir, report)
        report["slaytlar"].append(slide_info)
        report["toplam_metin"] += slide_info["metin_sayisi"]
        report["toplam_gorsel"] += slide_info["gorsel_sayisi"]

        print(f"   Metin blokları: {slide_info['metin_sayisi']}")
        print(f"   Görseller: {slide_info['gorsel_sayisi']}")

    # Markdown rapor oluştur
    _generate_report(report, output_dir)

    # JSON rapor oluştur (programatik erişim için)
    json_path = output_dir / "report.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)

    # Doğrulama
    _verify_output(report, output_dir)

    return report


# ─────────────────────────────────────────────
# Slayt İşleme
# ─────────────────────────────────────────────

def _process_slide(slide, slide_num, slide_dir, report):
    """Tek bir slaytı işler: metin ve görselleri ayırır."""
    slide_info = {
        "slayt_no": slide_num,
        "metin_sayisi": 0,
        "gorsel_sayisi": 0,
        "metin_bloklari": [],
        "gorsel_dosyalari": [],
    }

    text_parts = []
    image_count = 0

    # Şekilleri konum sırasına göre işle (üstten alta, soldan sağa)
    sorted_shapes = sorted(
        slide.shapes,
        key=lambda s: (s.top if s.top else 0, s.left if s.left else 0)
    )

    for shape_idx, shape in enumerate(sorted_shapes):
        try:
            # ── Grup şekilleri ──
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                g_texts, g_images, image_count = _process_group(
                    shape, slide_dir, image_count, report, slide_num
                )
                for gt in g_texts:
                    text_parts.append(gt)
                    slide_info["metin_bloklari"].append(gt)
                    slide_info["metin_sayisi"] += 1
                for gi in g_images:
                    slide_info["gorsel_dosyalari"].append(gi)
                    slide_info["gorsel_sayisi"] += 1
                continue

            # ── Görsel şekilleri ──
            if _is_picture(shape):
                image_count += 1
                img_filename = _save_image(shape, slide_dir, image_count)
                if img_filename:
                    slide_info["gorsel_dosyalari"].append(img_filename)
                    slide_info["gorsel_sayisi"] += 1
                    text_parts.append(f"![{img_filename}]({img_filename})")
                continue

            # ── Tablo şekilleri ──
            if shape.has_table:
                table_md = _extract_table(shape)
                if table_md.strip():
                    text_parts.append(table_md)
                    slide_info["metin_bloklari"].append(table_md)
                    slide_info["metin_sayisi"] += 1
                continue

            # ── Metin şekilleri ──
            if shape.has_text_frame:
                text_content = _extract_text(shape)
                if text_content.strip():
                    text_parts.append(text_content)
                    slide_info["metin_bloklari"].append(text_content)
                    slide_info["metin_sayisi"] += 1

        except Exception as e:
            report["hatalar"].append(
                f"Slayt {slide_num}, Şekil {shape_idx} ({shape.name}): {type(e).__name__} - {e}"
            )

    # Slayt içeriğini markdown dosyasına yaz
    content_path = slide_dir / "content.md"
    with open(content_path, "w", encoding="utf-8") as f:
        f.write(f"# Slayt {slide_num}\n\n")
        for part in text_parts:
            f.write(f"{part}\n\n")

    return slide_info


def _is_picture(shape):
    """Şeklin bir görsel olup olmadığını kontrol eder."""
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        try:
            _ = shape.image
            return True
        except (AttributeError, ValueError):
            return False
    # Genel kontrol
    try:
        _ = shape.image
        return True
    except (AttributeError, ValueError):
        return False


# ─────────────────────────────────────────────
# Metin Çıkarma
# ─────────────────────────────────────────────

def _extract_text(shape):
    """Bir şekilden metin çıkarır (Markdown formatlı)."""
    parts = []
    is_title = shape.name and any(
        t in shape.name.lower() for t in ("title", "başlık", "baslik")
    )

    for para in shape.text_frame.paragraphs:
        para_text = ""
        for run in para.runs:
            text = run.text
            if not text:
                continue
            # Kalın ve italik biçimlendirme
            if run.font.bold and run.font.italic:
                text = f"***{text}***"
            elif run.font.bold:
                text = f"**{text}**"
            elif run.font.italic:
                text = f"*{text}*"
            para_text += text

        # Eğer run'lardan metin gelmediyse doğrudan paragraf metnini al
        if not para_text and para.text:
            para_text = para.text

        if para_text.strip():
            # Başlık şekillerini markdown başlık olarak formatla
            if is_title and para.level == 0:
                para_text = f"## {para_text}"
            elif para.level and para.level > 0:
                indent = "  " * para.level
                para_text = f"{indent}- {para_text}"
            parts.append(para_text)

    return "\n".join(parts)


def _extract_table(shape):
    """Tabloyu Markdown formatında çıkarır."""
    table = shape.table
    rows = []
    for row_idx, row in enumerate(table.rows):
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip().replace("|", "\\|"))
        rows.append("| " + " | ".join(cells) + " |")
        if row_idx == 0:
            rows.append("| " + " | ".join(["---"] * len(row.cells)) + " |")
    return "\n".join(rows)


# ─────────────────────────────────────────────
# Görsel İşleme
# ─────────────────────────────────────────────

def _save_image(shape, slide_dir, image_count):
    """Görsel şeklini dosyaya kaydeder, bilgilerini döndürür."""
    try:
        image = shape.image
    except (AttributeError, ValueError):
        return None

    content_type = image.content_type or "image/png"
    ext = content_type.split("/")[-1]
    ext_map = {"jpeg": "jpg", "svg+xml": "svg", "x-emf": "emf", "x-wmf": "wmf"}
    ext = ext_map.get(ext, ext)

    filename = f"image_{image_count:02d}.{ext}"
    filepath = slide_dir / filename

    with open(filepath, "wb") as f:
        f.write(image.blob)

    size_kb = len(image.blob) / 1024
    detail = f"{size_kb:.1f} KB"

    if HAS_PIL and ext in ("png", "jpg", "gif", "bmp", "webp", "tiff"):
        try:
            img = Image.open(io.BytesIO(image.blob))
            detail = f"{img.width}x{img.height} px, {size_kb:.1f} KB"
        except Exception:
            pass

    print(f"   💾 Görsel kaydedildi: {filename} ({detail})")
    return filename


# ─────────────────────────────────────────────
# Grup Şekilleri
# ─────────────────────────────────────────────

def _process_group(group_shape, slide_dir, image_count, report, slide_num):
    """Grup şekillerini özyinelemeli olarak işler."""
    texts = []
    images = []

    for shape in group_shape.shapes:
        # İç içe grup
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            g_t, g_i, image_count = _process_group(
                shape, slide_dir, image_count, report, slide_num
            )
            texts.extend(g_t)
            images.extend(g_i)
            continue

        # Görsel
        if _is_picture(shape):
            try:
                image_count += 1
                img_filename = _save_image(shape, slide_dir, image_count)
                if img_filename:
                    images.append(img_filename)
            except Exception as e:
                report["hatalar"].append(
                    f"Slayt {slide_num}, Grup içi şekil: {e}"
                )
            continue

        # Metin
        if shape.has_text_frame:
            text = _extract_text(shape)
            if text.strip():
                texts.append(text)

        # Tablo
        if shape.has_table:
            table_md = _extract_table(shape)
            if table_md.strip():
                texts.append(table_md)

    return texts, images, image_count


# ─────────────────────────────────────────────
# Rapor Oluşturma
# ─────────────────────────────────────────────

def _generate_report(report, output_dir):
    """İşlem sonrası Markdown raporu oluşturur."""
    report_path = output_dir / "report.md"

    lines = [
        "# PPTX İşlem Raporu",
        "",
        f"**Dosya:** `{Path(report['dosya']).name}`  ",
        f"**Toplam Slayt:** {report['slayt_sayisi']}  ",
        f"**Slayt Boyutu:** {report['slayt_boyutu']}  ",
        f"**Toplam Metin Bloğu:** {report['toplam_metin']}  ",
        f"**Toplam Görsel:** {report['toplam_gorsel']}  ",
        "",
        "---",
        "",
    ]

    for slayt in report["slaytlar"]:
        sno = slayt["slayt_no"]
        lines.append(f"## Slayt {sno}")
        lines.append("")
        lines.append(f"- **Metin blokları:** {slayt['metin_sayisi']}")
        lines.append(f"- **Görseller:** {slayt['gorsel_sayisi']}")
        lines.append("")

        if slayt["metin_bloklari"]:
            lines.append("### Metin İçeriği")
            lines.append("")
            for txt in slayt["metin_bloklari"]:
                for line in txt.split("\n"):
                    lines.append(f"> {line}")
                lines.append("")

        if slayt["gorsel_dosyalari"]:
            lines.append("### Görseller")
            lines.append("")
            for img in slayt["gorsel_dosyalari"]:
                lines.append(f"![{img}](slide_{sno:02d}/{img})")
            lines.append("")

        lines.append("---")
        lines.append("")

    if report["hatalar"]:
        lines.append("## Hatalar")
        lines.append("")
        for hata in report["hatalar"]:
            lines.append(f"- {hata}")
        lines.append("")

    with open(report_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    print(f"\n📋 Rapor: {report_path}")


# ─────────────────────────────────────────────
# Doğrulama
# ─────────────────────────────────────────────

def _verify_output(report, output_dir):
    """Çıktı dosyalarını doğrular: her slayt için metin ve görsellerin varlığını kontrol eder."""
    print(f"\n{'═'*60}")
    print("DOĞRULAMA RAPORU")
    print(f"{'═'*60}")

    errors = []
    warnings = []
    ok_count = 0

    for slayt in report["slaytlar"]:
        sno = slayt["slayt_no"]
        slide_dir = output_dir / f"slide_{sno:02d}"
        slide_ok = True

        # Slayt klasörü var mı?
        if not slide_dir.exists():
            errors.append(f"Slayt {sno}: Klasör bulunamadı → {slide_dir}")
            continue

        # content.md var mı?
        content_file = slide_dir / "content.md"
        if not content_file.exists():
            errors.append(f"Slayt {sno}: content.md oluşturulamamış")
            slide_ok = False
        elif content_file.stat().st_size == 0:
            warnings.append(f"Slayt {sno}: content.md boş")

        # Boş slayt kontrolü
        if slayt["metin_sayisi"] == 0 and slayt["gorsel_sayisi"] == 0:
            warnings.append(f"Slayt {sno}: Boş slayt (metin ve görsel yok)")

        # Görsel dosyaları doğrula
        for img_file in slayt["gorsel_dosyalari"]:
            img_path = slide_dir / img_file
            if not img_path.exists():
                errors.append(f"Slayt {sno}: Görsel bulunamadı → {img_file}")
                slide_ok = False
            elif img_path.stat().st_size == 0:
                errors.append(f"Slayt {sno}: Görsel dosyası boş → {img_file}")
                slide_ok = False
            else:
                # Görsel bütünlüğünü kontrol et (PIL ile)
                if HAS_PIL and img_file.endswith((".png", ".jpg", ".gif", ".bmp")):
                    try:
                        img = Image.open(img_path)
                        img.verify()
                    except Exception as e:
                        errors.append(f"Slayt {sno}: Görsel bozuk → {img_file}: {e}")
                        slide_ok = False

        if slide_ok:
            ok_count += 1

    # Rapor dosyaları kontrolü
    for fname in ("report.md", "report.json"):
        fpath = output_dir / fname
        if not fpath.exists():
            errors.append(f"{fname} oluşturulamamış")

    # Sonuç yazdır
    print(f"\n  Toplam Slayt   : {report['slayt_sayisi']}")
    print(f"  Başarılı Slayt : {ok_count}")
    print(f"  Toplam Metin   : {report['toplam_metin']}")
    print(f"  Toplam Görsel  : {report['toplam_gorsel']}")

    if warnings:
        print(f"\n  ⚠️  UYARILAR ({len(warnings)}):")
        for w in warnings:
            print(f"    - {w}")

    if errors:
        print(f"\n  ❌ HATALAR ({len(errors)}):")
        for e in errors:
            print(f"    - {e}")
        print(f"\n  ❌ Doğrulama BAŞARISIZ! {len(errors)} hata bulundu.")
    else:
        print(f"\n  ✅ Doğrulama BAŞARILI! Tüm {report['slayt_sayisi']} slayt doğru işlendi.")

    print(f"\n  Çıktı klasörü: {output_dir}")
    print(f"{'═'*60}")

    return len(errors) == 0


# ─────────────────────────────────────────────
# CLI Giriş Noktası
# ─────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="PPTX dosyalarını analiz eder: metinleri ve görselleri ayırır.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Örnekler:
  python pptx_processor.py sunum.pptx
  python pptx_processor.py sunum.pptx --output-dir cikti/
  python pptx_processor.py sunum.pptx -o /tmp/sunum_analiz/
        """,
    )
    parser.add_argument("pptx_file", help="İşlenecek PPTX dosyasının yolu")
    parser.add_argument(
        "--output-dir", "-o",
        help="Çıktı klasörü (varsayılan: {dosya_adı}_output)",
        default=None,
    )

    args = parser.parse_args()
    report = process_pptx(args.pptx_file, args.output_dir)

    print(f"\nİşlem tamamlandı!")
    print(f"  Slayt: {report['slayt_sayisi']}  |  "
          f"Metin: {report['toplam_metin']}  |  "
          f"Görsel: {report['toplam_gorsel']}")
    if report["hatalar"]:
        print(f"  Hata: {len(report['hatalar'])}")
        sys.exit(1)


if __name__ == "__main__":
    main()
